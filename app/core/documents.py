"""문서 생성 — MD/TXT/코드, DOCX, XLSX, PPTX.

**입력 형식이 설계의 핵심이다.** 에이전트는 ReAct 액션 블록 안의 JSON 으로
인자를 넘기는데, 깊게 중첩된 구조(슬라이드 배열 안의 불릿 배열 …)는
따옴표·줄바꿈 이스케이프가 얽혀 자주 깨진다.

그래서 전부 **평평한 문자열**로 받는다:
  - DOCX/PPTX : 마크다운
  - XLSX      : CSV 텍스트

모델이 이미 잘 쓰는 형식이라 실패율이 낮고, 사람이 눈으로 검토하기도 쉽다.
"""
from __future__ import annotations

import csv
import io
import re
from typing import Any, Dict, List, Optional, Tuple

# 확장자별 MIME. 브라우저가 다운로드 시 올바른 앱으로 열게 한다.
MIME = {
    ".md": "text/markdown", ".txt": "text/plain", ".csv": "text/csv",
    ".json": "application/json", ".py": "text/x-python", ".sql": "text/plain",
    ".ts": "text/plain", ".js": "text/javascript", ".sh": "text/x-shellscript",
    ".yaml": "text/yaml", ".yml": "text/yaml", ".html": "text/html",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

_BAD_NAME = re.compile(r"[/\\\x00-\x1f]")
MAX_CHARS = 400_000


def clean_name(name: str, *, default: str, force_ext: Optional[str] = None) -> str:
    """파일명을 안전하게 만든다. 경로 구분자는 애초에 못 들어오게 한다."""
    name = _BAD_NAME.sub("_", (name or "").strip()) or default
    name = name.lstrip(".")[:120] or default
    if force_ext and not name.lower().endswith(force_ext):
        name = re.sub(r"\.[A-Za-z0-9]{1,5}$", "", name) + force_ext
    return name


def guess_mime(name: str) -> str:
    for ext, mime in MIME.items():
        if name.lower().endswith(ext):
            return mime
    return "text/plain"


# ---------------------------------------------------------------------------
# 마크다운 파싱 (공용)
# ---------------------------------------------------------------------------
_H = re.compile(r"^(#{1,6})\s+(.*)$")
_BULLET = re.compile(r"^(\s*)[-*+]\s+(.*)$")
_NUM = re.compile(r"^(\s*)\d+[.)]\s+(.*)$")
_TABLE_SEP = re.compile(r"^\s*\|?[\s:|-]+\|[\s:|-]*$")


def _inline(text: str) -> str:
    """굵게/기울임/코드 마크업만 벗겨낸다 (서식은 생성기가 따로 준다)."""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"\[(.+?)\]\((.+?)\)", r"\1 (\2)", text)
    return text.strip()


def _split_row(line: str) -> List[str]:
    cells = line.strip().strip("|").split("|")
    return [_inline(c.strip()) for c in cells]


def parse_blocks(md: str) -> List[Dict[str, Any]]:
    """마크다운을 블록 목록으로. heading / para / bullet / table / code"""
    blocks: List[Dict[str, Any]] = []
    lines = (md or "").replace("\r\n", "\n").split("\n")
    i = 0
    while i < len(lines):
        raw = lines[i]
        line = raw.rstrip()

        if line.strip().startswith("```"):
            lang = line.strip()[3:].strip()
            body = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith("```"):
                body.append(lines[i])
                i += 1
            i += 1
            blocks.append({"t": "code", "text": "\n".join(body), "lang": lang})
            continue

        h = _H.match(line)
        if h:
            blocks.append({"t": "heading", "level": len(h.group(1)),
                           "text": _inline(h.group(2))})
            i += 1
            continue

        # 표: 헤더 다음 줄이 구분선이어야 표로 본다
        if "|" in line and i + 1 < len(lines) and _TABLE_SEP.match(lines[i + 1]):
            header = _split_row(line)
            rows = []
            i += 2
            while i < len(lines) and "|" in lines[i] and lines[i].strip():
                rows.append(_split_row(lines[i]))
                i += 1
            blocks.append({"t": "table", "header": header, "rows": rows})
            continue

        b = _BULLET.match(line) or _NUM.match(line)
        if b:
            indent = len(b.group(1).replace("\t", "  ")) // 2
            blocks.append({"t": "bullet", "level": min(indent, 4),
                           "text": _inline(b.group(2))})
            i += 1
            continue

        if line.strip():
            blocks.append({"t": "para", "text": _inline(line.strip())})
        i += 1
    return blocks


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------
def build_docx(markdown: str, *, title: Optional[str] = None) -> bytes:
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    if title:
        doc.add_heading(title, level=0)

    for b in parse_blocks(markdown):
        kind = b["t"]
        if kind == "heading":
            doc.add_heading(b["text"], level=min(b["level"], 4))
        elif kind == "bullet":
            p = doc.add_paragraph(b["text"], style="List Bullet")
            if b["level"]:
                p.paragraph_format.left_indent = Pt(18 * (b["level"] + 1))
        elif kind == "code":
            p = doc.add_paragraph()
            run = p.add_run(b["text"])
            run.font.name = "Menlo"
            run.font.size = Pt(9)
        elif kind == "table" and b["header"]:
            t = doc.add_table(rows=1, cols=len(b["header"]))
            t.style = "Light Grid Accent 1"
            for c, name in zip(t.rows[0].cells, b["header"]):
                c.text = name
            for row in b["rows"]:
                cells = t.add_row().cells
                for c, v in zip(cells, row):
                    c.text = v
        else:
            doc.add_paragraph(b["text"])

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------
def _parse_csv(text: str) -> List[List[str]]:
    sample = (text or "")[:8192]
    try:
        delim = csv.Sniffer().sniff(sample, delimiters=",\t;|").delimiter
    except csv.Error:
        delim = ","
    return [r for r in csv.reader(io.StringIO(text or ""), delimiter=delim)]


def build_xlsx(csv_text: str, *, sheet_name: str = "Sheet1") -> Tuple[bytes, int, int]:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill
    from openpyxl.utils import get_column_letter

    rows = _parse_csv(csv_text)
    wb = Workbook()
    ws = wb.active
    ws.title = (sheet_name or "Sheet1")[:31]

    for r in rows:
        # 숫자로 보이는 값은 숫자로 넣는다 — 안 그러면 엑셀에서 합계가 안 된다
        ws.append([_coerce(v) for v in r])

    if rows:
        head = ws[1]
        fill = PatternFill("solid", fgColor="EEF1F6")
        for c in head:
            c.font = Font(bold=True)
            c.fill = fill
            c.alignment = Alignment(vertical="center")
        ws.freeze_panes = "A2"
        for idx in range(1, len(rows[0]) + 1):
            width = max((len(str(r[idx - 1])) for r in rows[:200]
                         if len(r) >= idx), default=8)
            ws.column_dimensions[get_column_letter(idx)].width = min(max(width + 2, 8), 48)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue(), len(rows), (len(rows[0]) if rows else 0)


def _coerce(v: str):
    s = (v or "").strip()
    if not s:
        return None
    try:
        if re.fullmatch(r"-?\d{1,3}(,\d{3})+(\.\d+)?|-?\d+(\.\d+)?", s):
            n = float(s.replace(",", ""))
            return int(n) if n.is_integer() and abs(n) < 2**53 else n
    except ValueError:
        pass
    return s


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------
def build_pptx(markdown: str, *, title: Optional[str] = None) -> Tuple[bytes, int]:
    """`#` 또는 `##` 가 새 슬라이드를 연다. 불릿은 그 슬라이드 본문이 된다."""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    blank = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[0]
    body_layout = prs.slide_layouts[1]

    blocks = parse_blocks(markdown)

    # 표지: 첫 블록이 h1 이면 그걸 쓰고, 아니면 title 인자를 쓴다
    cover = title
    if blocks and blocks[0]["t"] == "heading" and blocks[0]["level"] == 1:
        cover = cover or blocks[0]["text"]
    if cover:
        s = prs.slides.add_slide(title_layout)
        s.shapes.title.text = cover
        if len(s.placeholders) > 1:
            s.placeholders[1].text = ""

    slides = 1 if cover else 0
    body = None

    def new_slide(heading: str):
        nonlocal slides, body
        s = prs.slides.add_slide(body_layout)
        s.shapes.title.text = heading
        body = s.placeholders[1].text_frame
        body.clear()
        body._first = True          # 첫 문단은 이미 존재하므로 재사용
        slides += 1

    def add_line(text: str, level: int = 0, size: int = 16):
        nonlocal body
        if body is None:
            new_slide("")
        if getattr(body, "_first", False):
            p = body.paragraphs[0]
            body._first = False
        else:
            p = body.add_paragraph()
        p.text = text
        p.level = min(level, 4)
        for run in p.runs:
            run.font.size = Pt(size)

    skip_first_h1 = bool(cover and blocks and blocks[0]["t"] == "heading"
                         and blocks[0]["level"] == 1)
    for idx, b in enumerate(blocks):
        if skip_first_h1 and idx == 0:
            continue
        if b["t"] == "heading":
            new_slide(b["text"])
        elif b["t"] == "bullet":
            add_line(b["text"], b["level"])
        elif b["t"] == "table":
            add_line(" | ".join(b["header"]), 0, 14)
            for row in b["rows"][:8]:
                add_line(" | ".join(row), 1, 12)
        elif b["t"] == "code":
            for line in b["text"].split("\n")[:12]:
                add_line(line, 1, 11)
        else:
            add_line(b["text"], 0, 14)

    if slides == 0:
        prs.slides.add_slide(blank)
        slides = 1

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), slides
